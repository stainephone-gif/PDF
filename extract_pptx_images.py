#!/usr/bin/env python3
"""Extract all images from a PPTX file."""

import argparse
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def extract_images(pptx_path: str, output_dir: str | None = None) -> int:
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        print(f"Error: file not found: {pptx_path}", file=sys.stderr)
        return 1

    out_dir = Path(output_dir) if output_dir else pptx_path.parent / (pptx_path.stem + "_images")
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path)

    count = 0
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                filename = out_dir / f"slide{slide_num:03d}_img{count + 1:03d}.{ext}"
                filename.write_bytes(image.blob)
                print(f"Saved: {filename}")
                count += 1

    if count == 0:
        print("No images found in the presentation.")
    else:
        print(f"\nDone: {count} image(s) saved to '{out_dir}'")

    return 0


def main():
    parser = argparse.ArgumentParser(description="Extract images from a PPTX file")
    parser.add_argument("pptx", help="Path to the .pptx file")
    parser.add_argument(
        "-o", "--output",
        help="Output directory (default: <pptx_name>_images/ next to the file)",
    )
    args = parser.parse_args()
    sys.exit(extract_images(args.pptx, args.output))


if __name__ == "__main__":
    main()
