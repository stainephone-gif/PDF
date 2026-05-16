#!/usr/bin/env python3
"""Flask web UI for extracting images from PPTX files."""

import io
import zipfile
from pathlib import Path
from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # 100 MB


def extract_images_from_bytes(data: bytes) -> list[dict]:
    """Return list of {filename, blob} for every image in the PPTX."""
    prs = Presentation(io.BytesIO(data))
    images = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        img_idx = 1
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                images.append({
                    "filename": f"slide{slide_num:03d}_img{img_idx:03d}.{image.ext}",
                    "blob": image.blob,
                })
                img_idx += 1
    return images


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/extract", methods=["POST"])
def extract():
    if "file" not in request.files:
        return jsonify(error="Файл не выбран"), 400

    f = request.files["file"]
    if not f.filename or not f.filename.lower().endswith(".pptx"):
        return jsonify(error="Нужен файл с расширением .pptx"), 400

    try:
        images = extract_images_from_bytes(f.read())
    except Exception as exc:
        return jsonify(error=f"Ошибка при обработке файла: {exc}"), 500

    if not images:
        return jsonify(error="В презентации не найдено ни одной картинки"), 404

    # Pack into ZIP in memory
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for img in images:
            zf.writestr(img["filename"], img["blob"])
    buf.seek(0)

    stem = Path(f.filename).stem
    return send_file(
        buf,
        mimetype="application/zip",
        as_attachment=True,
        download_name=f"{stem}_images.zip",
    )


if __name__ == "__main__":
    app.run(debug=True, port=5000)
